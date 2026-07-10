"""
WRLD Presentation Builder - renders JSON specification to .pptx
Same architecture principle as DocumentBuilder: zero hardcoded content/
formatting decisions - everything comes from the spec, code just renders
it literally.
"""
import io
import base64
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
 
 
class PresentationBuilderError(Exception):
    """Raised when a slide spec is malformed (missing required fields, bad data)."""
 
 
class PresentationBuilder:
    # Standard python-pptx default template layout indices
    LAYOUT_TITLE = 0          # Title Slide (title + subtitle)
    LAYOUT_SECTION_HEADER = 2 # Section Header (title + subtitle, big divider look)
    LAYOUT_TITLE_CONTENT = 1  # Title and Content (title + bullet body placeholder)
    LAYOUT_BLANK = 6          # Blank
 
    def build(self, spec):
        prs = Presentation()
 
        # Widescreen 16:9 by default, matching modern proposal decks
        dims = spec.get('dimensions', {})
        prs.slide_width = Inches(dims.get('width', 13.333))
        prs.slide_height = Inches(dims.get('height', 7.5))
 
        for slide_spec in spec.get('slides', []):
            self._render_slide(prs, slide_spec)
 
        buf = io.BytesIO()
        prs.save(buf)
        buf.seek(0)
        return buf.getvalue()
 
    def _apply_text_format(self, text_frame, comp, default_size=18, default_bold=False):
        """Apply font size/color/bold/alignment from spec to every paragraph/run
        in a text frame - mirrors DocumentBuilder's per-component formatting."""
        font_size = comp.get('font_size', default_size)
        bold = comp.get('bold', default_bold)
        color_hex = comp.get('color')
        align = comp.get('alignment', 'left')
        align_map = {'left': PP_ALIGN.LEFT, 'center': PP_ALIGN.CENTER,
                     'right': PP_ALIGN.RIGHT, 'justify': PP_ALIGN.JUSTIFY}
 
        for para in text_frame.paragraphs:
            para.alignment = align_map.get(align, PP_ALIGN.LEFT)
            for run in para.runs:
                run.font.size = Pt(font_size)
                run.font.bold = bold
                if color_hex:
                    try:
                        r = int(color_hex[1:3], 16)
                        g = int(color_hex[3:5], 16)
                        b = int(color_hex[5:7], 16)
                        run.font.color.rgb = RGBColor(r, g, b)
                    except Exception:
                        pass
 
    def _set_hyperlink(self, text_frame, link):
        """Applies a hyperlink to every run in a text frame, if a link is given."""
        if not link:
            return
        for para in text_frame.paragraphs:
            for run in para.runs:
                run.hyperlink.address = link
 
    def _render_slide(self, prs, comp):
        slide_type = comp.get('type', 'bullet_content')
 
        if slide_type == 'title_slide':
            layout = prs.slide_layouts[self.LAYOUT_TITLE]
            slide = prs.slides.add_slide(layout)
            slide.shapes.title.text = comp.get('title', '')
            self._apply_text_format(slide.shapes.title.text_frame, comp.get('title_format', {}),
                                     default_size=40, default_bold=True)
            self._set_hyperlink(slide.shapes.title.text_frame, comp.get('title_link'))
            if len(slide.placeholders) > 1 and comp.get('subtitle'):
                sub_ph = slide.placeholders[1]
                sub_ph.text = comp.get('subtitle', '')
                self._apply_text_format(sub_ph.text_frame, comp.get('subtitle_format', {}),
                                         default_size=20, default_bold=False)
                self._set_hyperlink(sub_ph.text_frame, comp.get('subtitle_link'))
 
        elif slide_type == 'section_header':
            layout = prs.slide_layouts[self.LAYOUT_SECTION_HEADER]
            slide = prs.slides.add_slide(layout)
            slide.shapes.title.text = comp.get('title', '')
            self._apply_text_format(slide.shapes.title.text_frame, comp.get('title_format', {}),
                                     default_size=32, default_bold=True)
            if len(slide.placeholders) > 1 and comp.get('subtitle'):
                sub_ph = slide.placeholders[1]
                sub_ph.text = comp.get('subtitle', '')
                self._apply_text_format(sub_ph.text_frame, comp.get('subtitle_format', {}),
                                         default_size=16)
 
        elif slide_type == 'bullet_content':
            layout = prs.slide_layouts[self.LAYOUT_TITLE_CONTENT]
            slide = prs.slides.add_slide(layout)
            slide.shapes.title.text = comp.get('title', '')
            self._apply_text_format(slide.shapes.title.text_frame, comp.get('title_format', {}),
                                     default_size=28, default_bold=True)
 
            body = slide.placeholders[1]
            tf = body.text_frame
            tf.clear()
            bullets = comp.get('bullets', [])
            for i, bullet_item in enumerate(bullets):
                if isinstance(bullet_item, dict):
                    bullet_text = bullet_item.get('text', '')
                    bullet_level = bullet_item.get('level', comp.get('bullet_level', 0))
                    bullet_link = bullet_item.get('link')
                else:
                    bullet_text = bullet_item
                    bullet_level = comp.get('bullet_level', 0)
                    bullet_link = None
 
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                p.text = bullet_text
                p.level = bullet_level
                if bullet_link and p.runs:
                    p.runs[0].hyperlink.address = bullet_link
            self._apply_text_format(tf, comp.get('body_format', {}), default_size=18)
 
        elif slide_type == 'closing_slide':
            layout = prs.slide_layouts[self.LAYOUT_TITLE]
            slide = prs.slides.add_slide(layout)
            slide.shapes.title.text = comp.get('title', '')
            self._apply_text_format(slide.shapes.title.text_frame, comp.get('title_format', {}),
                                     default_size=36, default_bold=True)
            self._set_hyperlink(slide.shapes.title.text_frame, comp.get('title_link'))
            if len(slide.placeholders) > 1 and comp.get('subtitle'):
                sub_ph = slide.placeholders[1]
                sub_ph.text = comp.get('subtitle', '')
                self._apply_text_format(sub_ph.text_frame, comp.get('subtitle_format', {}),
                                         default_size=18)
                self._set_hyperlink(sub_ph.text_frame, comp.get('subtitle_link'))
 
        elif slide_type == 'image_slide':
            self._build_image_slide(prs, comp)
 
        elif slide_type == 'table_slide':
            self._build_table_slide(prs, comp)
 
        else:
            # Unknown slide type - blank slide with a text box noting it,
            # rather than silently dropping content
            layout = prs.slide_layouts[self.LAYOUT_BLANK]
            slide = prs.slides.add_slide(layout)
            box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(1))
            box.text_frame.text = f"[Unsupported slide type: {slide_type}]"
 
    def _build_image_slide(self, prs, comp):
        """
        Renders a title plus one or more images, laid out side by side.
        Each image entry must provide 'image_base64' (the calling agent is
        expected to have already fetched/prepared the image bytes - this
        builder does not fetch from URLs itself, to avoid adding a new
        outbound-request surface to a fixed-schema renderer).
        """
        layout = prs.slide_layouts[self.LAYOUT_BLANK]
        slide = prs.slides.add_slide(layout)
 
        w = prs.slide_width / 914400
        h = prs.slide_height / 914400
 
        title = comp.get('title', '')
        if title:
            title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(w - 1), Inches(1))
            title_box.text_frame.text = title
            self._apply_text_format(title_box.text_frame, comp.get('title_format', {}),
                                     default_size=28, default_bold=True)
            self._set_hyperlink(title_box.text_frame, comp.get('title_link'))
 
        images = comp.get('images', [])
        if not images:
            raise PresentationBuilderError("image_slide requires a non-empty 'images' list")
 
        top_in = 1.5
        side_margin = 0.5
        gap = 0.3
        available_width = w - (2 * side_margin)
        n = len(images)
        each_width = (available_width - gap * (n - 1)) / n
 
        for i, img in enumerate(images):
            b64 = img.get('image_base64')
            if not b64:
                raise PresentationBuilderError(f"images[{i}] requires 'image_base64'")
 
            try:
                img_bytes = base64.b64decode(b64)
            except Exception as e:
                raise PresentationBuilderError(f"images[{i}]: invalid base64 data ({e})")
 
            img_stream = io.BytesIO(img_bytes)
            left_in = side_margin + i * (each_width + gap)
            requested_width = min(img.get('width_in', each_width), each_width)
 
            pic = slide.shapes.add_picture(
                img_stream, Inches(left_in), Inches(top_in), width=Inches(requested_width)
            )
 
            caption = img.get('caption')
            if caption:
                cap_top_in = top_in + (pic.height / 914400) + 0.1
                cap_box = slide.shapes.add_textbox(
                    Inches(left_in), Inches(cap_top_in), Inches(each_width), Inches(0.4)
                )
                cap_box.text_frame.text = caption
                cap_box.text_frame.word_wrap = True
                self._apply_text_format(cap_box.text_frame, img.get('caption_format', {}),
                                         default_size=12)
 
    def _build_table_slide(self, prs, comp):
        """
        Renders a title plus a native pptx table. Column widths are only
        applied if explicitly provided in the spec - otherwise python-pptx's
        own even-distribution default is left alone.
        """
        layout = prs.slide_layouts[self.LAYOUT_BLANK]
        slide = prs.slides.add_slide(layout)
 
        w = prs.slide_width / 914400
        h = prs.slide_height / 914400
 
        title = comp.get('title', '')
        if title:
            title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(w - 1), Inches(1))
            title_box.text_frame.text = title
            self._apply_text_format(title_box.text_frame, comp.get('title_format', {}),
                                     default_size=28, default_bold=True)
            self._set_hyperlink(title_box.text_frame, comp.get('title_link'))
 
        rows = comp.get('rows', [])
        if not rows:
            raise PresentationBuilderError("table_slide requires a non-empty 'rows' list")
 
        n_rows = len(rows)
        n_cols = max(len(r) for r in rows)
        header_row = comp.get('header_row', True)
        font_size = comp.get('font_size', 14)
 
        top_in = 1.5
        table_width_in = w - 1.0
        table_height_in = min(h - top_in - 0.4, 0.5 * n_rows)
 
        gshape = slide.shapes.add_table(
            n_rows, n_cols,
            Inches(0.5), Inches(top_in),
            Inches(table_width_in), Inches(table_height_in),
        )
        table = gshape.table
 
        col_widths_in = comp.get('col_widths_in')
        if col_widths_in and len(col_widths_in) == n_cols:
            for ci, cw in enumerate(col_widths_in):
                table.columns[ci].width = Inches(cw)
 
        for ri, row_vals in enumerate(rows):
            for ci in range(n_cols):
                val = row_vals[ci] if ci < len(row_vals) else ''
                cell = table.cell(ri, ci)
                cell.text = str(val)
                for para in cell.text_frame.paragraphs:
                    for run in para.runs:
                        run.font.size = Pt(font_size)
                        if header_row and ri == 0:
                            run.font.bold = True
 
